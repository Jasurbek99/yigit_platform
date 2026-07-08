"""Build the YGT Platform 4-month progress presentation (.pptx).

Factual progress report: original plan -> time given -> what was delivered ->
scope added on top -> screen-by-screen walkthrough of the export process in the
operational (non-admin) views. Screenshot slides contain labeled placeholders
to paste real screenshots into.

Run:  python scripts/build_presentation.py
Out:  YGT_Progress_Report.pptx  (project root)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Palette ──────────────────────────────────────────────────────────────────
DARK   = RGBColor(0x1E, 0x2A, 0x38)   # slate (titles / bars)
ACCENT = RGBColor(0x2E, 0x7D, 0x32)   # green (growth / fills)
WARM   = RGBColor(0xE0, 0x7A, 0x3F)   # warm orange (highlights)
GRAY   = RGBColor(0x6B, 0x72, 0x80)   # secondary text
LIGHT  = RGBColor(0xF1, 0xF4, 0xF7)   # placeholder fill
TRACK  = RGBColor(0xE3, 0xE7, 0xEC)   # progress track
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


# ─── Helpers ────────────────────────────────────────────────────────────────
def _txt(slide, left, top, width, height, text, size, *, bold=False,
         color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return box


def _rect(slide, left, top, width, height, fill, *, line=None, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def header(slide, title, kicker=None):
    """Standard content-slide header: accent bar + title (+ optional kicker)."""
    _rect(slide, 0, 0, Inches(0.22), SH, ACCENT)
    if kicker:
        _txt(slide, Inches(0.6), Inches(0.35), Inches(11), Inches(0.4),
             kicker.upper(), 12, bold=True, color=WARM)
        _txt(slide, Inches(0.6), Inches(0.72), Inches(12), Inches(0.9),
             title, 30, bold=True, color=DARK)
    else:
        _txt(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.9),
             title, 30, bold=True, color=DARK)


def bullets(slide, items, *, top=Inches(2.0), left=Inches(0.8),
            width=Inches(11.7), size=18, gap=0.62):
    """items: list of (text, level) or plain strings (level 0)."""
    box = slide.shapes.add_textbox(left, top, width, SH - top - Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        text, level = (item if isinstance(item, tuple) else (item, 0))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap * 18)
        p.level = level
        bullet = "•   " if level == 0 else "–   "
        r = p.add_run()
        r.text = bullet + text
        f = r.font
        f.name = FONT
        f.size = Pt(size if level == 0 else size - 2)
        f.bold = level == 0 and item == items[0]  # no-op safety
        f.color.rgb = DARK if level == 0 else GRAY
    return box


def new_slide():
    return prs.slides.add_slide(BLANK)


def notes(slide, text):
    """Attach speaker notes — the talking points for that slide."""
    slide.notes_slide.notes_text_frame.text = text


# ─── 1. Title ─────────────────────────────────────────────────────────────────
s = new_slide()
_rect(s, 0, 0, SW, SH, DARK)
_rect(s, 0, Inches(5.05), SW, Inches(0.08), ACCENT)
_txt(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(0.5),
     "YGT HOLDING  ·  DIGITAL TRANSFORMATION", 16, bold=True, color=WARM)
_txt(s, Inches(0.9), Inches(2.45), Inches(11.5), Inches(1.6),
     "Export Platform — Progress Report", 46, bold=True, color=WHITE)
_txt(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(0.8),
     "What was planned · time given · what was delivered in 4 months",
     20, color=RGBColor(0xC8, 0xD0, 0xD8))
_txt(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.5),
     "Built solo with AI assistance   ·   February – June 2026", 15,
     color=RGBColor(0x9A, 0xA5, 0xB1))
notes(s, "Thank you for the time. This is a short, honest report of the export "
         "platform: what we set out to build, how much time we had, what is "
         "done after four months, and what I added along the way. I'll walk "
         "through the actual export process at the end so you can see it the "
         "way the team uses it every day.")


# ─── 2. The original plan ──────────────────────────────────────────────────────
s = new_slide()
header(s, "The Original Plan", "Roadmap v1.0 — February 2026")
bullets(s, [
    "One unified Django + React platform to replace Excel, phone calls and WhatsApp across the whole value chain.",
    "Five interconnected projects:",
    ("P1 Greenhouse  ·  P2 Transport  ·  P3 Export  ·  P4 Contracts  ·  P5 Finance", 1),
    "Planned timeline: 12 months.",
    "Planned team (per roadmap): 2–3 developers + 1 data analyst + 0.5 QA.",
    "Defined priority order: P3 Export first — highest operational pain, most data available, every other project depends on it.",
], top=Inches(2.0))
notes(s, "The original roadmap is one platform covering the whole value chain "
         "in five projects, planned over twelve months for a team of two to "
         "three developers plus an analyst. The plan itself said to build "
         "Export first, because it is the biggest daily pain and everything "
         "else depends on its data. I followed that priority.")


# ─── 3. Time & resources: plan vs reality ──────────────────────────────────────
s = new_slide()
header(s, "Time & Resources — Plan vs. Reality")

rows, cols = 4, 3
tbl_w = Inches(12.0)
tbl = s.shapes.add_table(rows, cols, Inches(0.7), Inches(1.9), tbl_w, Inches(2.4)).table
tbl.columns[0].width = Inches(3.4)
tbl.columns[1].width = Inches(4.3)
tbl.columns[2].width = Inches(4.3)
data = [
    ("", "Planned (roadmap)", "Actual"),
    ("Team", "2–3 developers + analyst + QA", "1 developer + AI"),
    ("Timeline", "12 months", "~4 months so far"),
    ("Discovery / planning", "Interviews + analysis up front", "Compressed — coding began mid-March"),
]
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = data[r][c]
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0] if para.runs else para.add_run()
        run.font.name = FONT
        run.font.size = Pt(15)
        if r == 0:
            run.font.bold = True
            run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = DARK
        else:
            run.font.bold = (c == 0)
            run.font.color.rgb = DARK if c != 2 else ACCENT
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT

_txt(s, Inches(0.7), Inches(4.7), Inches(12), Inches(0.45),
     "Timeline", 15, bold=True, color=WARM)
bullets(s, [
    "Late February 2026 — planning started.",
    "Mid-March 2026 — coding started (the learn-the-business / discovery window was very short).",
    "End May – early June 2026 — first results reviewed.  →  roughly 3.5–4 months of actual build.",
], top=Inches(5.1), size=15, gap=0.5)
notes(s, "Two numbers to keep in mind for the rest of this. First, the team: "
         "the plan assumed two to three developers and an analyst — this was "
         "built by one person with AI assistance. Second, the time: planning "
         "started late February, coding mid-March, and we are reviewing at the "
         "end of May. That is under four months of build, with a very short "
         "window up front to learn the business in detail. I mention this not "
         "as an excuse — the result stands on its own — but so the scope is "
         "judged against the time and the people that were actually behind it.")


# ─── 4. Delivered in 4 months (headline) ────────────────────────────────────────
s = new_slide()
_rect(s, 0, 0, SW, SH, DARK)
_rect(s, 0, 0, Inches(0.22), SH, ACCENT)
_txt(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
     "DELIVERED IN ~4 MONTHS", 14, bold=True, color=WARM)
# Big number
_txt(s, Inches(0.8), Inches(1.7), Inches(6), Inches(2.0),
     "~40%", 110, bold=True, color=WHITE)
_txt(s, Inches(0.8), Inches(3.9), Inches(11.5), Inches(0.7),
     "of the full 12-month, 5-project scope", 24,
     color=RGBColor(0xC8, 0xD0, 0xD8))
bullets_box = s.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.3))
tf = bullets_box.text_frame; tf.word_wrap = True
for i, line in enumerate([
    "The highest-priority project — P3 Export — is ~90% complete and in daily production use.",
    "Delivered solo + AI what the plan budgeted for a 2–3 person team over a longer horizon.",
    "Value is front-loaded exactly as the roadmap intended: highest-pain project first.",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(10)
    r = p.add_run(); r.text = "•   " + line
    r.font.name = FONT; r.font.size = Pt(17); r.font.color.rgb = RGBColor(0xDD, 0xE3, 0xE9)
notes(s, "The headline: in under four months, roughly forty percent of a "
         "twelve-month, five-project plan is done. But it is not a flat forty "
         "percent everywhere — the work was concentrated where it matters most. "
         "The Export project, the one the roadmap said to do first, is about "
         "ninety percent complete and already running in daily use. So the most "
         "important and most-used part of the whole plan is essentially live.")


# ─── 5. Completion by project (progress bars) ───────────────────────────────────
s = new_slide()
header(s, "Completion by Project")
projects = [
    ("P3  Export", 90, "Core platform — in production, used daily by 9 roles"),
    ("P5  Finance / Executive", 40, "Executive dashboards built; ERP / 1C integration pending"),
    ("P4  Contracts", 35, "Contracts + invoices + debt view; document auto-gen pending"),
    ("P1  Greenhouse", 30, "Harvest planning + blocks; mobile PWA + analytics pending"),
    ("P2  Transport", 5, "Not started — fleet registry / GPS"),
]
top = Inches(2.05)
row_h = Inches(0.95)
label_w = Inches(3.1)
bar_left = Inches(3.8)
bar_w = Inches(6.0)
bar_h = Inches(0.34)
for i, (name, pct, note) in enumerate(projects):
    y = Emu(int(top) + int(row_h) * i)
    _txt(s, Inches(0.7), Emu(int(y)), label_w, Inches(0.4), name, 16, bold=True)
    # track
    by = Emu(int(y) + Inches(0.05))
    _rect(s, bar_left, by, bar_w, bar_h, TRACK, rounded=True)
    fill_w = Emu(int(int(bar_w) * pct / 100))
    color = ACCENT if pct >= 50 else (WARM if pct >= 25 else RGBColor(0xC0, 0x39, 0x2B))
    if pct > 0:
        _rect(s, bar_left, by, fill_w, bar_h, color, rounded=True)
    _txt(s, Emu(int(bar_left) + int(bar_w) + int(Inches(0.15))), Emu(int(y)),
         Inches(0.9), Inches(0.4), f"{pct}%", 16, bold=True, color=color)
    _txt(s, bar_left, Emu(int(y) + int(Inches(0.46))), Inches(8.5), Inches(0.4),
         note, 11.5, color=GRAY)
notes(s, "Breaking the forty percent down by project. Export is at ninety "
         "percent and in production. The executive and finance dashboards are "
         "about forty percent — the screens exist, the accounting-system "
         "integrations are still ahead. Contracts is about thirty-five percent "
         "— contracts, invoices and the debt view are in, the document "
         "auto-generation is the next big piece. Greenhouse is early, and "
         "Transport has not started. This is deliberate: highest-value first, "
         "exactly as the plan ordered it.")


# ─── 6. Scope added (not planned) ───────────────────────────────────────────────
s = new_slide()
header(s, "Scope Added During the Project", "Not in the original plan")
bullets(s, [
    "Kanban boards + Tasks — requested by management mid-project.",
    "Comments, @mentions and per-cell discussion — a full team-collaboration layer.",
    "\"Sheet\" — a Google-Sheets-style grid, built to win real user adoption (next slide).",
    "\"My Tasks\" personal board for each role's daily work.",
    "Feedback module + Worklog / time tracking.",
    "Configurable role & field permissions + delegated user management.",
    "Three-language interface (Turkmen / Russian / English).",
    "Error tracking (Sentry) and a full audit log.",
], top=Inches(1.95), size=16, gap=0.5)
notes(s, "A large part of the four months went into things that were not in "
         "the original plan. Some were asked for during the project — the "
         "kanban boards and tasks. Others came from talking to the people who "
         "actually use the system every day. The biggest of these is the "
         "spreadsheet view, which I'll explain on the next slide.")


# ─── 7. Why the Sheet (adoption) ────────────────────────────────────────────────
s = new_slide()
header(s, "Why the Spreadsheet View", "Solving the #1 project risk: user adoption")
bullets(s, [
    "The roadmap flagged \"user resistance\" as the highest-likelihood risk.",
    "Employees were clear: they would not move off their familiar spreadsheets.",
    "Solution: build a spreadsheet-style interface — the platform feels familiar, but the data is now connected, validated and real-time.",
    "Same way of working the team already knows — with the control a single source of truth gives.",
    "This is an adoption strategy, not a rebuild of Excel: behind the grid sits the full lifecycle, permissions, quota and reporting engine.",
], top=Inches(2.0), size=17)
notes(s, "I want to be honest about the spreadsheet view, because at first "
         "glance it can look like we just rebuilt Excel. We did not. The "
         "biggest risk to this whole project was that people would refuse to "
         "leave their spreadsheets — the roadmap itself listed this as the "
         "number-one risk. The team told me directly they would not switch. So "
         "I made the platform feel like what they already know, while "
         "underneath it is connected, validated, and a single source of truth. "
         "This is how we get real adoption instead of people quietly going back "
         "to their own files.")


# ─── 8. Section divider — the export process ────────────────────────────────────
s = new_slide()
_rect(s, 0, 0, SW, SH, DARK)
_rect(s, Inches(0.8), Inches(2.7), Inches(2.4), Inches(0.1), ACCENT)
_txt(s, Inches(0.8), Inches(3.0), Inches(11.5), Inches(1.4),
     "The Export Process, End to End", 40, bold=True, color=WHITE)
_txt(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.0),
     "The operational role views — what the 9 daily users actually see "
     "(distinct from the admin configuration view).", 18,
     color=RGBColor(0xC8, 0xD0, 0xD8), italic=True)
notes(s, "Now I'll show the actual export process, screen by screen, the way "
         "the team sees it when they log in — not the admin view. This is a "
         "single shipment travelling from harvest, through documents and "
         "customs, to sale and final report. Each screen here is a real, "
         "working part of that flow.")


# ─── 9+. Showcase slides with screenshot placeholders ───────────────────────────
    # (title, description, capture-instruction: which screen + which login)
showcase = [
    ("Dashboard", "Daily landing page for every role — key stats, alerts (missing reports, quota, document queue), active shipments and routes.",
     "SCREEN: Home / Dashboard (right after login)  ·  login as export_manager"),
    ("Daily Harvest & Draft Creation", "Supply side: harvest entered per block, then turned into one-truck shipment drafts — the start of the export pipeline.",
     "SCREEN: Daily Harvest board / Draft pool  ·  login as loading_dept_head or export_manager"),
    ("The Sheet — Main Working Surface", "Spreadsheet-style grid all 9 roles work in. Inline editing, frozen rows/columns, copy/paste, undo, per-cell comments — but every change is validated and shared in real time.",
     "SCREEN: Shipments → Sheet view  ·  login as export_manager (NOT admin)"),
    ("Kanban Board — Phase Tracking", "Each shipment flows through 7 phases (Plan → Prep → Docs → Load → Transit → Dest → Close), with average time-in-phase per column.",
     "SCREEN: Shipments → Board (Kanban)  ·  login as export_manager"),
    ("Shipment Detail — 13-Step Lifecycle", "Full record of one shipment with the 13-step timeline, weights, firm splits, documents and status history.",
     "SCREEN: open any shipment → Detail (show the 13-step timeline)  ·  login as export_manager"),
    ("Quota Management", "Government 1:10 quota tracking per firm — usage, remaining balance, per-firm progress and alerts.",
     "SCREEN: Quota dashboard  ·  login as export_manager"),
    ("Weekly Harvest Planning", "Block × day planning grid — plan vs. actual, feeding the export pipeline.",
     "SCREEN: Weekly harvest planning grid  ·  login as export_manager"),
    ("My Tasks", "Each employee's personal board: to-do / in progress / blocked / done-today, drawn from tasks assigned across shipments.",
     "SCREEN: My Tasks board  ·  login as any role that has tasks"),
    ("Comments & Collaboration", "Threaded comments with @user / @role mentions attached to specific cells — replaces the WhatsApp coordination.",
     "SCREEN: open the Comments drawer on a Sheet cell  ·  login as export_manager"),
    ("Executive Dashboard", "Management view: revenue, outstanding debt by firm, firm-risk matrix, route profitability and quota status.",
     "SCREEN: Boss / executive dashboard  ·  login as director or boss"),
    ("Contracts & Invoices", "Contract registry and invoice tracking with the outstanding-debt breakdown.",
     "SCREEN: Contracts list + Invoices  ·  login as export_manager or finansist"),
    ("Admin & Permissions", "Configurable role/page/field permissions, reference data, seasons and audit log — the administration layer behind the operational views.",
     "SCREEN: Admin → Permissions / Users  ·  login as admin"),
]

for title, desc, capture in showcase:
    s = new_slide()
    header(s, title)
    # description
    _txt(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.9), desc, 15, color=GRAY)
    # placeholder
    ph_top = Inches(2.55)
    ph = _rect(s, Inches(0.9), ph_top, Inches(11.5), Inches(4.4), LIGHT,
               line=TRACK, rounded=True)
    _txt(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(0.7),
         "PUT SCREENSHOT HERE", 22, bold=True, color=DARK,
         align=PP_ALIGN.CENTER)
    _txt(s, Inches(1.2), Inches(4.75), Inches(10.9), Inches(0.9),
         capture, 14, bold=True, color=WARM, align=PP_ALIGN.CENTER)
    notes(s, desc + "  (Capture instruction on the slide: " + capture + ")")


# ─── What's left ────────────────────────────────────────────────────────────────
s = new_slide()
header(s, "What's Left on the Roadmap", "Honest remaining scope")
bullets(s, [
    "P4 — Document auto-generation engine (CMR, invoice, customs, fito, CT-1, TIR): the daily 2–3 hours of manual filling.",
    "P2 — Transport & Fleet: truck/driver registry, availability board, GPS tracking.",
    "P1 — Greenhouse depth: mobile/offline harvest input, yield analytics.",
    "P5 — Finance integrations: Logo Tiger ERP (read) and 1C credit data.",
    "P3 — Polish and full-season hardening as real volume hits the system.",
], top=Inches(2.0), size=17)
notes(s, "I want to be clear about what is not done yet, so there are no "
         "surprises. The biggest remaining piece is the document "
         "auto-generation — the two to three hours of manual filling every "
         "day. Then transport and fleet, greenhouse depth, and the accounting "
         "integrations. The work is well-understood; it is a question of time "
         "and priority from here.")


# ─── Summary ────────────────────────────────────────────────────────────────────
s = new_slide()
_rect(s, 0, 0, SW, SH, DARK)
_rect(s, 0, Inches(1.4), SW, Inches(0.08), ACCENT)
_txt(s, Inches(0.8), Inches(0.55), Inches(11.5), Inches(0.8),
     "Summary", 36, bold=True, color=WHITE)
box = s.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.7), Inches(5.0))
tf = box.text_frame; tf.word_wrap = True
for i, line in enumerate([
    "In ~4 months, solo + AI: ~40% of a 12-month, 5-project plan.",
    "The core export platform (P3) is ~90% done and in daily production use.",
    "Significant unplanned scope added on top — kanban, tasks, collaboration, the spreadsheet view — much of it to drive real adoption.",
    "The full export process is live in the operational role views, end to end.",
    "Remaining work is well-scoped: documents, transport, greenhouse depth, finance integrations.",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(14)
    r = p.add_run(); r.text = "•   " + line
    r.font.name = FONT; r.font.size = Pt(18); r.font.color.rgb = RGBColor(0xDD, 0xE3, 0xE9)
notes(s, "To summarise: in under four months, one person with AI built about "
         "forty percent of a twelve-month, five-project plan — with the core "
         "export platform live and in daily use, plus a lot of extra work to "
         "make sure people actually adopt it. The full export process is "
         "working end to end in the operational views I just showed. The "
         "remaining work is clear and scoped. Thank you — I'm happy to take "
         "questions or walk through any screen in more detail.")


out = "YGT_Progress_Report.pptx"
prs.save(out)
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides")
